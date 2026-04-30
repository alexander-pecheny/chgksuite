from pathlib import Path

import toml
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.util import Inches as PptxInches
from pptx.util import Pt as PptxPt

from chgksuite.common import DefaultArgs
from chgksuite.composer.pptx import PptxExporter


ROOT = Path(__file__).resolve().parents[1]
RESOURCES = ROOT / "chgksuite" / "resources"


def _merge_config(base, updates):
    for key, value in updates.items():
        if value is None:
            base.pop(key, None)
            continue
        if isinstance(value, dict):
            base.setdefault(key, {})
            _merge_config(base[key], value)
        else:
            base[key] = value
    return base


def _config_path(tmp_path, updates=None):
    if not updates:
        return str(RESOURCES / "pptx_config.toml")
    config = toml.load(RESOURCES / "pptx_config.toml")
    config["template_path"] = str(RESOURCES / "template.pptx")
    _merge_config(config, updates)
    path = tmp_path / "pptx_config.toml"
    path.write_text(toml.dumps(config), encoding="utf8")
    return str(path)


def _pptx_args(tmp_path, font=None, config_updates=None):
    return DefaultArgs(
        pptx_config=_config_path(tmp_path, config_updates),
        labels_file=str(RESOURCES / "labels_ru.toml"),
        regexes_file=str(RESOURCES / "regexes_ru.json"),
        language="ru",
        replace_no_break_spaces="on",
        replace_no_break_hyphens="on",
        font=font,
    )


def _export_pptx(tmp_path, structure, font=None, config_updates=None):
    outfilename = tmp_path / "out.pptx"
    exporter = PptxExporter(
        structure,
        _pptx_args(tmp_path, font=font, config_updates=config_updates),
        {"tmp_dir": str(tmp_path), "targetdir": str(tmp_path)},
    )
    exporter.export(str(outfilename))
    return Presentation(str(outfilename))


def test_title_slide_uses_full_height_centered_textbox(tmp_path):
    prs = _export_pptx(
        tmp_path,
        [
            (
                "heading",
                "Альфа. «Зеркало» первого игрового дня основной дисциплины СтудЧР-2026",
            ),
            ("Question", {"question": "Вопрос.", "answer": "Ответ."}),
        ],
    )

    title = prs.slides[0].shapes.title

    assert title.text_frame.auto_size == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    assert title.text_frame.vertical_anchor == MSO_VERTICAL_ANCHOR.MIDDLE
    assert round(title.top / 914400, 2) == 0.8
    assert round(title.height / 914400, 2) == 6.1
    assert round(title.width / 914400, 2) > 0
    assert [run.font.size.pt for run in title.text_frame.paragraphs[0].runs] == [60.0]
    assert not any(shape.name.startswith("Subtitle") for shape in prs.slides[0].shapes)


def test_pptx_textboxes_shrink_text_and_stamp_run_sizes(tmp_path):
    prs = _export_pptx(
        tmp_path,
        [
            ("heading", "Тестовый пакет"),
            ("section", "Тур 2"),
            (
                "editor",
                "Редакторы: Александр Сновский, Артём Сапожников, "
                "Александр Зайцев, Ольга Шиншинова, Мария Аристова.",
            ),
            (
                "meta",
                "Редакторы благодарят за тестирование вопросов и очень ценные "
                "замечания: Андрея Багдуева, Дмитрия Батова, Марию Богуш, "
                "Виталия Буковского, Михаила Гриценко, Дарью Жукову, "
                "Фёдора Журавлёва, Григория Зырянова, Алексея Ковбу, "
                "Евгению Колпащикову, Юстину Кустовскую, Маргариту Лузину, "
                "Михаила Малкина, Веру Монину, Ерлана Мухамеджанова, "
                "Илью Орлова, Никиту Пеговса, Кирилла Платонова, "
                "Юрия Разумова, Дмитрия Селянина, Дмитрия Слободянюка, "
                "Наиля Фарукшина, Ксению и Эдуарда Шагалов и Максима Шиловского.",
            ),
            ("Question", {"question": "Вопрос.", "answer": "Ответ."}),
        ],
    )

    textboxes = [
        shape
        for shape in prs.slides[1].shapes
        if hasattr(shape, "text_frame") and "Редакторы:" in shape.text
    ]
    assert len(textboxes) == 1
    textbox = textboxes[0]

    assert textbox.text_frame.auto_size == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    assert "<a:normAutofit" in textbox.element.txBody.xml
    assert "<a:spAutoFit" not in textbox.element.txBody.xml

    run_sizes = [
        run.font.size.pt
        for paragraph in textbox.text_frame.paragraphs
        for run in paragraph.runs
        if run.text.strip()
    ]
    assert run_sizes
    assert set(run_sizes) == {32.0}


def test_pptx_font_override_replaces_config_font(tmp_path):
    prs = _export_pptx(
        tmp_path,
        [
            ("heading", "Тестовый пакет"),
            ("Question", {"question": "Вопрос.", "answer": "Ответ."}),
        ],
        font="Times New Roman",
    )

    title = prs.slides[0].shapes.title
    assert title.text_frame.paragraphs[0].runs[0].font.name == "Times New Roman"

    question_runs = [
        run
        for shape in prs.slides[1].shapes
        if hasattr(shape, "text_frame") and "Вопрос." in shape.text
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
        if run.text.strip()
    ]
    assert question_runs
    assert {run.font.name for run in question_runs} == {"Times New Roman"}


def test_douplet_list_gets_extra_break_before_numbering(tmp_path):
    prs = _export_pptx(
        tmp_path,
        [
            ("heading", "Тестовый пакет"),
            (
                "Question",
                {
                    "question": [
                        "Дуплет. Два вопроса по 30 секунд каждый.",
                        [
                            "Первый подвопрос.",
                            "Второй подвопрос.",
                        ],
                    ],
                    "answer": ["Первый.", "Второй."],
                },
            ),
        ],
    )

    question_text = next(
        shape.text
        for shape in prs.slides[1].shapes
        if hasattr(shape, "text_frame") and "Дуплет." in shape.text
    )

    assert "каждый.\n\n1. Первый подвопрос." in question_text


def test_douplet_list_numbering_style_is_configurable(tmp_path):
    prs = _export_pptx(
        tmp_path,
        [
            ("heading", "Тестовый пакет"),
            (
                "Question",
                {
                    "question": [
                        "Блиц.",
                        [
                            "Первый подвопрос.",
                            "Второй подвопрос.",
                        ],
                    ],
                    "answer": ["Первый.", "Второй."],
                },
            ),
        ],
        config_updates={"list": {"numbering_style": "a)"}},
    )

    question_text = next(
        shape.text
        for shape in prs.slides[2].shapes
        if hasattr(shape, "text_frame") and "Блиц." in shape.text
    )

    assert "Блиц.\n\na) Первый подвопрос.\nb) Второй подвопрос." in question_text


def test_handout_slide_uses_handout_config_and_hides_label_by_default(tmp_path):
    prs = _export_pptx(
        tmp_path,
        [
            ("heading", "Тестовый пакет"),
            (
                "Question",
                {
                    "question": "[Раздаточный материал: Текст раздатки]\nВопрос?",
                    "answer": "Ответ.",
                },
            ),
        ],
        config_updates={"handout": {"font_size": 24, "align": "center"}},
    )

    handout_shape = next(
        shape
        for shape in prs.slides[1].shapes
        if hasattr(shape, "text_frame") and "Текст раздатки" in shape.text
    )
    paragraph = handout_shape.text_frame.paragraphs[0]

    assert handout_shape.text == "Текст раздатки"
    assert paragraph.alignment == PP_ALIGN.CENTER
    assert {run.font.size.pt for run in paragraph.runs if run.text.strip()} == {24.0}


def test_handout_label_can_be_included_from_config(tmp_path):
    prs = _export_pptx(
        tmp_path,
        [
            ("heading", "Тестовый пакет"),
            (
                "Question",
                {
                    "question": "[Раздаточный материал: Текст раздатки]\nВопрос?",
                    "answer": "Ответ.",
                },
            ),
        ],
        config_updates={"handout": {"include_label": True}},
    )

    handout_text = next(
        shape.text
        for shape in prs.slides[1].shapes
        if hasattr(shape, "text_frame") and "Текст раздатки" in shape.text
    )

    assert handout_text == "[Раздаточный материал: Текст раздатки]"


def test_inline_handout_uses_handout_config_when_not_separate_slide(tmp_path):
    prs = _export_pptx(
        tmp_path,
        [
            ("heading", "Тестовый пакет"),
            (
                "Question",
                {
                    "question": (
                        "[Раздаточный материал:\n"
                        "Klemperer\n"
                        "]\n"
                        "Чтобы спастись от гестапо, супруги Клемперер решили "
                        "подделать свои документы."
                    ),
                    "answer": "Ответ.",
                },
            ),
        ],
        config_updates={
            "add_handout_on_separate_slide": False,
            "force_text_size_question": 24,
            "font": {"default_size": None},
            "handout": {"include_label": False, "font_size": 32, "align": "center"},
        },
    )

    question_shape = next(
        shape
        for shape in prs.slides[1].shapes
        if hasattr(shape, "text_frame") and "Klemperer" in shape.text
    )
    handout_paragraph = question_shape.text_frame.paragraphs[0]
    question_paragraph = question_shape.text_frame.paragraphs[1]

    assert handout_paragraph.text == "Klemperer"
    assert handout_paragraph.alignment == PP_ALIGN.CENTER
    assert {
        run.font.size.pt for run in handout_paragraph.runs if run.text.strip()
    } == {32.0}
    assert handout_paragraph.space_after.pt == 18.0
    assert question_paragraph.text.startswith("Чтобы спастись")
    assert {
        run.font.size.pt for run in question_paragraph.runs if run.text.strip()
    } == {24.0}


def test_handout_image_scale_and_spacing_apply_to_question_slide(tmp_path):
    image_path = ROOT / "tests" / "test.jpg"
    prs = _export_pptx(
        tmp_path,
        [
            ("heading", "Тестовый пакет"),
            (
                "Question",
                {
                    "question": (
                        f"[Раздаточный материал: (img w=3in {image_path})]\n"
                        "Перед вами картинка. Назовите ее."
                    ),
                    "answer": "Ответ.",
                },
            ),
        ],
        config_updates={"handout": {"image_scale": 1.5, "space_after": 24}},
    )

    question_slide = next(
        slide
        for slide in prs.slides
        if any(
            hasattr(shape, "text_frame")
            and "Перед вами картинка" in shape.text.replace("\xa0", " ")
            for shape in slide.shapes
        )
    )
    picture = next(shape for shape in question_slide.shapes if shape.shape_type == 13)
    question_shape = next(
        shape
        for shape in question_slide.shapes
        if hasattr(shape, "text_frame")
        and "Перед вами картинка" in shape.text.replace("\xa0", " ")
    )

    assert picture.width == PptxInches(4.5)
    assert question_shape.top == picture.top + picture.height + PptxPt(24)


def test_legacy_pptx_config_sizes_and_disabled_autolayout_do_not_overlap(tmp_path):
    image_path = ROOT / "tests" / "test.jpg"
    prs = _export_pptx(
        tmp_path,
        [
            ("heading", "Тестовый пакет"),
            (
                "Question",
                {
                    "question": (
                        f"[Раздаточный материал: (img {image_path})]\n"
                        "Перед вами картинка. Назовите ее."
                    ),
                    "answer": "Ответ.",
                },
            ),
        ],
        config_updates={
            "add_plug": False,
            "add_handout_on_separate_slide": False,
            "disable_autolayout": True,
            "force_text_size_question": 24,
            "force_text_size_answer": 20,
            "text_size_grid": {"default": 24},
            "number_textbox": {"font_size": 28},
            "font": {"default_size": None},
        },
    )

    question_slide = prs.slides[1]
    picture = next(shape for shape in question_slide.shapes if shape.shape_type == 13)
    question_shape = next(
        shape
        for shape in question_slide.shapes
        if hasattr(shape, "text_frame")
        and "Перед вами картинка" in shape.text.replace("\xa0", " ")
    )
    number_shape = next(
        shape
        for shape in question_slide.shapes
        if hasattr(shape, "text_frame") and shape.text == "1"
    )

    assert question_shape.top >= picture.top + picture.height
    assert {
        run.font.size.pt
        for paragraph in question_shape.text_frame.paragraphs
        for run in paragraph.runs
        if run.text.strip()
    } == {24.0}
    assert number_shape.text_frame.paragraphs[0].runs[0].font.size.pt == 28.0

    answer_shape = next(
        shape
        for shape in prs.slides[2].shapes
        if hasattr(shape, "text_frame") and "Ответ." in shape.text
    )
    assert {
        run.font.size.pt
        for paragraph in answer_shape.text_frame.paragraphs
        for run in paragraph.runs
        if run.text.strip()
    } == {20.0}
